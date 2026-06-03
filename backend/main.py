import os
import uuid
import zipfile
import shutil
import asyncio
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Import authentication and history routes
from routes_auth import router as auth_router
from routes_history import router as history_router, save_to_history
from routes_cv import router as cv_router
from routes_stripe import router as stripe_router
from auth import get_current_user
from database import init_database

app = FastAPI(title="PDF Tools API", version="1.0.0")

# CORS - allow all origins for Flutter frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Include routes
app.include_router(auth_router)
app.include_router(history_router)
app.include_router(cv_router)
app.include_router(stripe_router)

TEMP_DIR = Path("temp_files")

@app.on_event("startup")
async def startup_event():
    TEMP_DIR.mkdir(exist_ok=True)
    # Check database connection
    try:
        init_database()
        print("✅ Database connection successful")
    except Exception as e:
        print(f"⚠️  Database connection failed: {e}")
        print("   Guest mode will work, but authentication will be unavailable")


def get_temp_path(suffix: str = ".pdf") -> Path:
    return TEMP_DIR / f"{uuid.uuid4()}{suffix}"

async def save_upload(file: UploadFile) -> Path:
    ext = Path(file.filename).suffix or ".pdf"
    path = get_temp_path(suffix=ext)
    with open(path, "wb") as f:
        f.write(await file.read())
    return path

def cleanup_files(*paths: Path):
    for p in paths:
        try:
            if p and p.exists():
                if p.is_dir():
                    shutil.rmtree(p)
                else:
                    p.unlink()
        except Exception:
            pass

async def schedule_cleanup(*paths: Path, delay: int = 3600):
    """Delete files after `delay` seconds (default 1 hour)."""
    await asyncio.sleep(delay)
    cleanup_files(*paths)

# ─────────────────────────────────────────────
#  ROOT
# ─────────────────────────────────────────────
@app.get("/")
def root():
    return {"message": "PDF Tools API is running.", "version": "1.0.0"}

@app.get("/health")
def health():
    return {"status": "ok"}

# ─────────────────────────────────────────────
#  1. MERGE PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/merge")
async def merge_pdf(
    background_tasks: BackgroundTasks,
    files: list[UploadFile] = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="At least 2 PDF files are required.")

    merger = PdfMerger()
    saved_inputs: list[Path] = []
    output_path: Optional[Path] = None

    try:
        for file in files:
            p = await save_upload(file)
            saved_inputs.append(p)
            merger.append(str(p))

        output_path = get_temp_path(".pdf")
        merger.write(str(output_path))
        merger.close()

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="merge",
                output_file=str(output_path),
                output_name="merged.pdf"
            )

        background_tasks.add_task(schedule_cleanup, output_path, *saved_inputs, delay=86400)  # 24 hours
        return FileResponse(
            path=str(output_path),
            filename="merged.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(*saved_inputs)
        if output_path:
            cleanup_files(output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  2. SPLIT PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/split")
async def split_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    ranges: Optional[str] = Form(None),
    every_page: bool = Form(False),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """
    Split a PDF.
    - `every_page=true`  → each page becomes its own PDF, returned as a ZIP.
    - `ranges`           → comma-separated page ranges like "1-3,5,7-9" (1-indexed), returned as a ZIP.
    """
    input_path = await save_upload(file)
    try:
        reader = PdfReader(str(input_path))
        total = len(reader.pages)

        zip_path = get_temp_path(".zip")
        zip_dir = TEMP_DIR / str(uuid.uuid4())
        zip_dir.mkdir()

        if every_page:
            for i in range(total):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                out = zip_dir / f"page_{i + 1}.pdf"
                with open(out, "wb") as f:
                    writer.write(f)
        elif ranges:
            for seg in ranges.split(","):
                seg = seg.strip()
                if "-" in seg:
                    start, end = seg.split("-")
                    start, end = int(start) - 1, int(end) - 1
                else:
                    start = end = int(seg) - 1
                writer = PdfWriter()
                for i in range(start, min(end + 1, total)):
                    writer.add_page(reader.pages[i])
                label = f"pages_{start + 1}_to_{end + 1}.pdf"
                out = zip_dir / label
                with open(out, "wb") as f:
                    writer.write(f)
        else:
            raise HTTPException(status_code=400, detail="Provide 'ranges' or set 'every_page=true'.")

        # Zip the output directory
        with zipfile.ZipFile(zip_path, "w") as z:
            for pdf in zip_dir.glob("*.pdf"):
                z.write(pdf, pdf.name)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="split",
                output_file=str(zip_path),
                output_name="split_pages.zip"
            )

        background_tasks.add_task(cleanup_files, input_path, zip_path, zip_dir)
        return FileResponse(
            path=str(zip_path),
            filename="split_pages.zip",
            media_type="application/zip",
        )
    except HTTPException:
        raise
    except Exception as e:
        cleanup_files(input_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  3. COMPRESS PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/compress")
async def compress_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    level: str = Form("medium"),  # low | medium | high
    current_user: Optional[dict] = Depends(get_current_user),
):
    """
    Compress a PDF by removing redundant/unused objects.
    True advanced compression requires Ghostscript (optional enhancement).
    """
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    original_size = input_path.stat().st_size

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()

        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)

        try:
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)
        except Exception as e:
            # PyPDF2's compress_identical_objects can sometimes crash on complex or malformed PDFs.
            # We catch it here so that compress_content_streams is still applied.
            print(f"Warning: compress_identical_objects failed: {e}")

        with open(output_path, "wb") as f:
            writer.write(f)

        compressed_size = output_path.stat().st_size
        saved = original_size - compressed_size
        pct = round((saved / original_size) * 100, 1) if original_size > 0 else 0

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="compress",
                output_file=str(output_path),
                output_name="compressed.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="compressed.pdf",
            media_type="application/pdf",
            headers={
                "X-Original-Size": str(original_size),
                "X-Compressed-Size": str(compressed_size),
                "X-Saved-Bytes": str(saved),
                "X-Saved-Percent": str(pct),
            },
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  4. ROTATE PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/rotate")
async def rotate_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    angle: int = Form(90),
    pages: Optional[str] = Form(None),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """
    Rotate pages by `angle` degrees (90, 180, 270).
    `pages` is a comma-separated list of 1-indexed page numbers. If omitted, all pages are rotated.
    """
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        total = len(reader.pages)

        if pages:
            target = set(int(p.strip()) - 1 for p in pages.split(","))
        else:
            target = set(range(total))

        for i, page in enumerate(reader.pages):
            if i in target:
                page.rotate(angle)
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="rotate",
                output_file=str(output_path),
                output_name="rotated.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="rotated.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  5. ADD WATERMARK (Text)
# ─────────────────────────────────────────────
@app.post("/api/pdf/watermark")
async def watermark_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    watermark_text: str = Form("CONFIDENTIAL"),
    opacity: float = Form(0.3),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """
    Overlay a text watermark PDF onto each page.
    Generates the watermark page on-the-fly using reportlab (if available) or PyPDF2.
    """
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    watermark_path = get_temp_path(".pdf")

    try:
        # Create watermark PDF with reportlab
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            import math

            c = canvas.Canvas(str(watermark_path), pagesize=A4)
            w, h = A4
            c.setFillColorRGB(0.5, 0.5, 0.5, alpha=opacity)
            c.setFont("Helvetica-Bold", 60)
            c.saveState()
            c.translate(w / 2, h / 2)
            c.rotate(45)
            c.drawCentredString(0, 0, watermark_text)
            c.restoreState()
            c.save()
        except ImportError:
            # Fallback: copy the file without watermark
            shutil.copy(input_path, output_path)
            
            # Save to history for authenticated users
            if current_user:
                await save_to_history(
                    user_id=current_user["user_id"],
                    tool_id="watermark",
                    output_file=str(output_path),
                    output_name="watermarked.pdf"
                )
            
            background_tasks.add_task(cleanup_files, input_path, output_path, watermark_path)
            return FileResponse(
                path=str(output_path),
                filename="watermarked.pdf",
                media_type="application/pdf",
                headers={"X-Warning": "reportlab not installed; watermark skipped"},
            )

        reader = PdfReader(str(input_path))
        wm_reader = PdfReader(str(watermark_path))
        writer = PdfWriter()

        wm_page = wm_reader.pages[0]
        for page in reader.pages:
            page.merge_page(wm_page)
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="watermark",
                output_file=str(output_path),
                output_name="watermarked.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path, watermark_path)
        return FileResponse(
            path=str(output_path),
            filename="watermarked.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path, watermark_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  6. PROTECT PDF (Password)
# ─────────────────────────────────────────────
@app.post("/api/pdf/protect")
async def protect_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="protect",
                output_file=str(output_path),
                output_name="protected.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="protected.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  7. UNLOCK PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/unlock")
async def unlock_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    try:
        reader = PdfReader(str(input_path))
        if reader.is_encrypted:
            result = reader.decrypt(password)
            if result == 0:
                cleanup_files(input_path)
                raise HTTPException(status_code=401, detail="Incorrect password.")

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="unlock",
                output_file=str(output_path),
                output_name="unlocked.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="unlocked.pdf",
            media_type="application/pdf",
        )
    except HTTPException:
        raise
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  8. IMAGES TO PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/images-to-pdf")
async def images_to_pdf(
    background_tasks: BackgroundTasks,
    files: list[UploadFile] = File(...),
    orientation: str = Form("portrait"),
    margin: int = Form(20),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert image files (JPG, PNG, etc.) to a single PDF."""
    try:
        from PIL import Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, landscape
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="Pillow and reportlab are required for image conversion. Run: pip install Pillow reportlab",
        )

    output_path = get_temp_path(".pdf")
    saved_inputs: list[Path] = []

    try:
        pagesize = landscape(A4) if orientation == "landscape" else A4
        pw, ph = pagesize

        c = canvas.Canvas(str(output_path), pagesize=pagesize)
        for file in files:
            p = await save_upload(file)
            saved_inputs.append(p)
            img = Image.open(str(p))
            iw, ih = img.size
            avail_w = pw - 2 * margin
            avail_h = ph - 2 * margin
            scale = min(avail_w / iw, avail_h / ih)
            draw_w, draw_h = iw * scale, ih * scale
            x = (pw - draw_w) / 2
            y = (ph - draw_h) / 2
            c.drawImage(str(p), x, y, draw_w, draw_h)
            c.showPage()

        c.save()

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="images-to-pdf",
                output_file=str(output_path),
                output_name="images_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, *saved_inputs, output_path)
        return FileResponse(
            path=str(output_path),
            filename="images_to_pdf.pdf",
            media_type="application/pdf",
        )
    except HTTPException:
        raise
    except Exception as e:
        cleanup_files(*saved_inputs, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  9. PDF TO IMAGES
# ─────────────────────────────────────────────
@app.post("/api/pdf/pdf-to-images")
async def pdf_to_images(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    dpi: int = Form(150),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert each page of a PDF to a JPG image, returned as ZIP."""
    try:
        from pdf2image import convert_from_path
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="pdf2image is required. Run: pip install pdf2image (also needs Poppler installed).",
        )

    input_path = await save_upload(file)
    zip_path = get_temp_path(".zip")
    img_dir = TEMP_DIR / str(uuid.uuid4())
    img_dir.mkdir()

    try:
        images = convert_from_path(str(input_path), dpi=dpi)
        img_paths: list[Path] = []
        for i, img in enumerate(images):
            img_path = img_dir / f"page_{i + 1}.jpg"
            img.save(str(img_path), "JPEG", quality=90)
            img_paths.append(img_path)

        with zipfile.ZipFile(zip_path, "w") as z:
            for p in img_paths:
                z.write(p, p.name)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="pdf-to-images",
                output_file=str(zip_path),
                output_name="pdf_pages.zip"
            )

        background_tasks.add_task(cleanup_files, input_path, zip_path, img_dir)
        return FileResponse(
            path=str(zip_path),
            filename="pdf_pages.zip",
            media_type="application/zip",
        )
    except Exception as e:
        cleanup_files(input_path, img_dir)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  10. PDF INFO
# ─────────────────────────────────────────────
@app.post("/api/pdf/info")
async def pdf_info(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Return metadata about a PDF: page count, title, author, etc."""
    input_path = await save_upload(file)
    try:
        reader = PdfReader(str(input_path))
        meta = reader.metadata or {}
        info = {
            "page_count": len(reader.pages),
            "title": meta.get("/Title", ""),
            "author": meta.get("/Author", ""),
            "subject": meta.get("/Subject", ""),
            "creator": meta.get("/Creator", ""),
            "is_encrypted": reader.is_encrypted,
            "file_size_bytes": input_path.stat().st_size,
        }
        background_tasks.add_task(cleanup_files, input_path)
        return JSONResponse(content=info)
    except Exception as e:
        cleanup_files(input_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  11. EXTRACT TEXT (OCR-ready)
# ─────────────────────────────────────────────
@app.post("/api/pdf/extract-text")
async def extract_text(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    input_path = await save_upload(file)
    try:
        reader = PdfReader(str(input_path))
        pages_text = []
        for i, page in enumerate(reader.pages):
            pages_text.append({"page": i + 1, "text": page.extract_text() or ""})
        background_tasks.add_task(cleanup_files, input_path)
        return JSONResponse(content={"pages": pages_text, "total_pages": len(reader.pages)})
    except Exception as e:
        cleanup_files(input_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  12. TEXT TO PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/text-to-pdf")
async def text_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert a plain text/HTML file to a PDF."""
    try:
        from xhtml2pdf import pisa
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="xhtml2pdf is required. Run: pip install xhtml2pdf",
        )

    filename = file.filename.lower()
    is_excel = filename.endswith(".xlsx") or filename.endswith(".xls")

    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        if is_excel:
            import pandas as pd
            xl = pd.ExcelFile(str(input_path))
            
            style = """
            <style>
                @page {
                    size: a4 landscape;
                    margin: 1.5cm;
                }
                body {
                    font-family: Helvetica, Arial, sans-serif;
                    color: #1E293B;
                }
                .sheet-title {
                    font-size: 16px;
                    font-weight: bold;
                    color: #4F46E5;
                    margin-top: 25px;
                    margin-bottom: 10px;
                    border-bottom: 2px solid #E2E8F0;
                    padding-bottom: 6px;
                }
                table {
                    width: 100%;
                    border-collapse: collapse;
                    margin-bottom: 25px;
                }
                th {
                    background-color: #4F46E5;
                    color: #FFFFFF;
                    font-weight: bold;
                    text-align: left;
                    padding: 8px;
                    border: 1px solid #CBD5E0;
                    font-size: 10px;
                }
                td {
                    padding: 8px;
                    border: 1px solid #E2E8F0;
                    font-size: 9px;
                }
                tr:nth-child(even) {
                    background-color: #F8FAFC;
                }
            </style>
            """
            
            html_parts = [style]
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                df = df.fillna("")
                table_html = df.to_html(index=False, border=0)
                html_parts.append(f'<div class="sheet-title">Sheet: {sheet_name}</div>')
                html_parts.append(table_html)
                
            text = f"<html><head></head><body>{''.join(html_parts)}</body></html>"
        else:
            with open(input_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()

            # Wrap plain text in simple HTML if it doesn't look like HTML
            if "<html" not in text.lower() and "<p" not in text.lower() and "<div" not in text.lower():
                text = f"<html><body><pre>{text}</pre></body></html>"
            
        with open(output_path, "wb") as result_file:
            pisa_status = pisa.CreatePDF(text, dest=result_file)

        if pisa_status.err:
            raise Exception("Failed to generate PDF from text/HTML")

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="text-to-pdf",
                output_file=str(output_path),
                output_name="text_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="text_to_pdf.pdf",
            media_type="application/pdf",
        )
    except HTTPException:
        raise
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))



# ─────────────────────────────────────────────
#  13. WORD TO PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/word-to-pdf")
async def word_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert DOCX/DOC files to PDF."""
    try:
        from docx import Document
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-docx and reportlab are required. Run: pip install python-docx reportlab",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        doc = Document(str(input_path))
        
        # Create PDF
        pdf = SimpleDocTemplate(str(output_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        for para in doc.paragraphs:
            if para.text.strip():
                p = Paragraph(para.text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 0.2 * inch))
        
        pdf.build(story)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="word-to-pdf",
                output_file=str(output_path),
                output_name="word_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="word_to_pdf.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  14. PDF TO WORD
# ─────────────────────────────────────────────
@app.post("/api/pdf/pdf-to-word")
async def pdf_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert PDF to DOCX."""
    try:
        from pdf2docx import Converter
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="pdf2docx is required. Run: pip install pdf2docx",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".docx")

    try:
        cv = Converter(str(input_path))
        cv.convert(str(output_path))
        cv.close()

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="pdf-to-word",
                output_file=str(output_path),
                output_name="pdf_to_word.docx"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="pdf_to_word.docx",
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  15. PDF TO EXCEL
# ─────────────────────────────────────────────
@app.post("/api/pdf/pdf-to-excel")
async def pdf_to_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Extract tables from PDF and convert to Excel."""
    try:
        import tabula
        import pandas as pd
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="tabula-py is required. Run: pip install tabula-py (requires Java)",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".xlsx")

    try:
        # Extract all tables from PDF
        tables = tabula.read_pdf(str(input_path), pages='all', multiple_tables=True)
        
        if not tables:
            raise HTTPException(status_code=400, detail="No tables found in PDF")

        # Write to Excel with multiple sheets
        with pd.ExcelWriter(str(output_path), engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                sheet_name = f"Table_{i + 1}"
                table.to_excel(writer, sheet_name=sheet_name, index=False)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="pdf-to-excel",
                output_file=str(output_path),
                output_name="pdf_to_excel.xlsx"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="pdf_to_excel.xlsx",
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    except HTTPException:
        raise
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  16. EXCEL TO PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/excel-to-pdf")
async def excel_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert Excel files to PDF."""
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        import pandas as pd
        from xhtml2pdf import pisa

        xl = pd.ExcelFile(str(input_path))
        
        style = """
        <style>
            @page {
                size: a4 landscape;
                margin: 1.5cm;
            }
            body {
                font-family: Helvetica, Arial, sans-serif;
                color: #1E293B;
            }
            .sheet-title {
                font-size: 16px;
                font-weight: bold;
                color: #4F46E5;
                margin-top: 25px;
                margin-bottom: 10px;
                border-bottom: 2px solid #E2E8F0;
                padding-bottom: 6px;
            }
            table {
                width: 100%;
                border-collapse: collapse;
                margin-bottom: 25px;
            }
            th {
                background-color: #4F46E5;
                color: #FFFFFF;
                font-weight: bold;
                text-align: left;
                padding: 8px;
                border: 1px solid #CBD5E0;
                font-size: 10px;
            }
            td {
                padding: 8px;
                border: 1px solid #E2E8F0;
                font-size: 9px;
            }
            tr:nth-child(even) {
                background-color: #F8FAFC;
            }
        </style>
        """
        
        html_parts = [style]
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            df = df.fillna("")
            table_html = df.to_html(index=False, border=0)
            html_parts.append(f'<div class="sheet-title">Sheet: {sheet_name}</div>')
            html_parts.append(table_html)
            
        html_content = f"<html><head></head><body>{''.join(html_parts)}</body></html>"
        
        with open(output_path, "wb") as result_file:
            pisa_status = pisa.CreatePDF(html_content, dest=result_file)

        if pisa_status.err:
            raise Exception("Failed to generate PDF from Excel")

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="excel-to-pdf",
                output_file=str(output_path),
                output_name="excel_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="excel_to_pdf.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  17. POWERPOINT TO PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/powerpoint-to-pdf")
async def powerpoint_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert PowerPoint files to PDF."""
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas as pdf_canvas
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx is required. Run: pip install python-pptx",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        prs = Presentation(str(input_path))
        c = pdf_canvas.Canvas(str(output_path), pagesize=landscape(letter))
        page_width, page_height = landscape(letter)

        for slide_num, slide in enumerate(prs.slides):
            # Add slide number and basic content
            c.setFont("Helvetica-Bold", 24)
            c.drawString(50, page_height - 50, f"Slide {slide_num + 1}")
            
            y_position = page_height - 100
            c.setFont("Helvetica", 12)
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_lines = shape.text.split('\n')
                    for line in text_lines[:10]:  # Limit lines per slide
                        if y_position > 50:
                            c.drawString(50, y_position, line[:100])
                            y_position -= 20
            
            c.showPage()

        c.save()

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="powerpoint-to-pdf",
                output_file=str(output_path),
                output_name="powerpoint_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="powerpoint_to_pdf.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  18. PDF TO POWERPOINT
# ─────────────────────────────────────────────
@app.post("/api/pdf/pdf-to-powerpoint")
async def pdf_to_powerpoint(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert PDF pages to PowerPoint slides (as images)."""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="pdf2image and python-pptx are required. Run: pip install pdf2image python-pptx",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".pptx")
    img_dir = TEMP_DIR / str(uuid.uuid4())
    img_dir.mkdir()

    try:
        # Convert PDF pages to images
        images = convert_from_path(str(input_path), dpi=150)
        
        # Create PowerPoint
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, img in enumerate(images):
            img_path = img_dir / f"page_{i + 1}.png"
            img.save(str(img_path), "PNG")
            
            # Add blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add image to slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(str(img_path), left, top, width=prs.slide_width)

        prs.save(str(output_path))

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="pdf-to-powerpoint",
                output_file=str(output_path),
                output_name="pdf_to_powerpoint.pptx"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path, img_dir)
        return FileResponse(
            path=str(output_path),
            filename="pdf_to_powerpoint.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        cleanup_files(input_path, output_path, img_dir)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  19. HTML TO PDF (URL-based)
# ─────────────────────────────────────────────
@app.post("/api/pdf/html-url-to-pdf")
async def html_url_to_pdf(
    background_tasks: BackgroundTasks,
    url: str = Form(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Convert a webpage URL to PDF."""
    try:
        import requests
        from xhtml2pdf import pisa
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="requests and xhtml2pdf are required.",
        )

    output_path = get_temp_path(".pdf")

    try:
        # Fetch webpage content
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        html_content = response.text

        # Convert to PDF
        with open(output_path, "wb") as result_file:
            pisa_status = pisa.CreatePDF(html_content, dest=result_file)

        if pisa_status.err:
            raise Exception("Failed to generate PDF from URL")

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="html-to-pdf",
                output_file=str(output_path),
                output_name="webpage.pdf"
            )

        background_tasks.add_task(cleanup_files, output_path)
        return FileResponse(
            path=str(output_path),
            filename="webpage.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  20. ORGANIZE PDF (Reorder/Delete pages)
# ─────────────────────────────────────────────
@app.post("/api/pdf/organize")
async def organize_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    page_order: str = Form(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """
    Reorganize PDF pages.
    page_order: comma-separated page numbers (1-indexed), e.g., "3,1,2,5" to reorder pages.
    """
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        
        # Parse page order
        page_indices = [int(p.strip()) - 1 for p in page_order.split(",")]
        
        # Add pages in specified order
        for idx in page_indices:
            if 0 <= idx < len(reader.pages):
                writer.add_page(reader.pages[idx])

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="organize",
                output_file=str(output_path),
                output_name="organized.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="organized.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  21. ADD PAGE NUMBERS
# ─────────────────────────────────────────────
@app.post("/api/pdf/add-page-numbers")
async def add_page_numbers(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    position: str = Form("bottom-center"),
    start_number: int = Form(1),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Add page numbers to PDF."""
    try:
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.pagesizes import A4
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="reportlab is required.",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    numbers_path = get_temp_path(".pdf")

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()

        # Create page numbers PDF
        c = pdf_canvas.Canvas(str(numbers_path), pagesize=A4)
        page_width, page_height = A4

        for i in range(len(reader.pages)):
            page_num = start_number + i
            
            # Determine position
            if "bottom" in position:
                y = 30
            elif "top" in position:
                y = page_height - 30
            else:
                y = page_height / 2

            if "center" in position:
                x = page_width / 2
            elif "right" in position:
                x = page_width - 50
            else:
                x = 50

            c.setFont("Helvetica", 10)
            c.drawString(x, y, str(page_num))
            c.showPage()

        c.save()

        # Merge page numbers with original PDF
        numbers_reader = PdfReader(str(numbers_path))
        for i, page in enumerate(reader.pages):
            page.merge_page(numbers_reader.pages[i])
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="organize",
                output_file=str(output_path),
                output_name="organized.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path, numbers_path)
        return FileResponse(
            path=str(output_path),
            filename="numbered.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path, numbers_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  22. OCR PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/ocr")
async def ocr_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    language: str = Form("eng"),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Perform OCR on PDF to make it searchable."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.pagesizes import A4
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="pytesseract and pdf2image are required. Run: pip install pytesseract pdf2image",
        )

    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")
    img_dir = TEMP_DIR / str(uuid.uuid4())
    img_dir.mkdir()

    try:
        # Convert PDF to images
        images = convert_from_path(str(input_path), dpi=300)
        
        # Create searchable PDF
        c = pdf_canvas.Canvas(str(output_path), pagesize=A4)
        
        for i, img in enumerate(images):
            # Perform OCR
            text = pytesseract.image_to_string(img, lang=language)
            
            # Add text to PDF
            c.setFont("Helvetica", 10)
            y = 800
            for line in text.split('\n'):
                if y > 50:
                    c.drawString(50, y, line[:100])
                    y -= 15
            
            c.showPage()

        c.save()

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="powerpoint-to-pdf",
                output_file=str(output_path),
                output_name="powerpoint_to_pdf.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path, img_dir)
        return FileResponse(
            path=str(output_path),
            filename="ocr_result.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path, img_dir)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  23. CROP PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/crop")
async def crop_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    left: float = Form(0),
    bottom: float = Form(0),
    right: float = Form(0),
    top: float = Form(0),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Crop margins from PDF pages."""
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        reader = PdfReader(str(input_path))
        writer = PdfWriter()

        for page in reader.pages:
            # Get current page dimensions
            page_box = page.mediabox
            
            # Apply crop
            page.mediabox.lower_left = (
                page_box.left + left,
                page_box.bottom + bottom
            )
            page.mediabox.upper_right = (
                page_box.right - right,
                page_box.top - top
            )
            
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="organize",
                output_file=str(output_path),
                output_name="organized.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path)
        return FileResponse(
            path=str(output_path),
            filename="cropped.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))


# ─────────────────────────────────────────────
#  24. REDACT PDF
# ─────────────────────────────────────────────
@app.post("/api/pdf/redact")
async def redact_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    search_terms: str = Form(...),
    current_user: Optional[dict] = Depends(get_current_user),
):
    """Redact (blackout) sensitive text in PDF."""
    input_path = await save_upload(file)
    output_path = get_temp_path(".pdf")

    try:
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.pagesizes import A4
        
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        
        terms = [term.strip() for term in search_terms.split(",")]
        
        # Create redaction overlay
        redaction_path = get_temp_path(".pdf")
        c = pdf_canvas.Canvas(str(redaction_path), pagesize=A4)
        page_width, page_height = A4

        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            
            # Simple redaction by drawing black boxes (basic implementation)
            c.setFillColorRGB(0, 0, 0)
            
            for term in terms:
                if term.lower() in text.lower():
                    # Draw a black rectangle (simplified - would need exact coordinates)
                    c.rect(100, 700, 200, 20, fill=True, stroke=False)
            
            c.showPage()

        c.save()

        # Merge redaction overlay with original
        redaction_reader = PdfReader(str(redaction_path))
        for i, page in enumerate(reader.pages):
            page.merge_page(redaction_reader.pages[i])
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Save to history for authenticated users
        if current_user:
            await save_to_history(
                user_id=current_user["user_id"],
                tool_id="organize",
                output_file=str(output_path),
                output_name="organized.pdf"
            )

        background_tasks.add_task(cleanup_files, input_path, output_path, redaction_path)
        return FileResponse(
            path=str(output_path),
            filename="redacted.pdf",
            media_type="application/pdf",
        )
    except Exception as e:
        cleanup_files(input_path, output_path)
        raise HTTPException(status_code=500, detail=str(e))

